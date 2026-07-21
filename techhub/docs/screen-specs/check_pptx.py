# -*- coding: utf-8 -*-
"""자가 점검 v2:
(1) 텍스트 프레임 경계 초과, (2) 슬라이드 경계 이탈,
(3) 본문 와이어프레임 내 문장형 텍스트(>4자) 잔존 = 무언어 규칙 위반,
(4) 모든 본문 슬라이드에 화면 전체 레이아웃(브라우저 프레임) 포함 여부."""
import sys, math, unicodedata
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

EMU = 914400
SW, SH = 13.333, 7.5
MARGIN = 0.42
WIRE_R = MARGIN + 5.65     # 와이어프레임 우측 경계
DESC_X = MARGIN + 5.65 + 0.3


def is_wide(ch):
    return unicodedata.east_asian_width(ch) in ("W", "F", "A")


def tw_pt(text, size):
    return sum(size * (1.0 if is_wide(ch) else 0.55) for ch in text)


def emu_in(v):
    try:
        return v / EMU
    except Exception:
        return 0.05


def check(path):
    prs = Presentation(path)
    issues = []
    for idx, slide in enumerate(prs.slides, 1):
        has_table = any(sh.has_table for sh in slide.shapes)
        is_body = has_table
        has_frame = False
        for sh in slide.shapes:
            try:
                x = sh.left / EMU; y = sh.top / EMU; w = sh.width / EMU; h = sh.height / EMU
            except TypeError:
                continue
            # (4) 브라우저 프레임 감지
            if abs(x - MARGIN) < 0.15 and 5.3 < w < 5.9 and h > 3.0:
                has_frame = True
            # (2) 경계 이탈
            if x < -0.03 or y < -0.03 or x + w > SW + 0.03 or y + h > SH + 0.03:
                issues.append(f"[S{idx}] OUT-OF-BOUNDS '{_lbl(sh)}' r={x+w:.2f} b={y+h:.2f}")
            if sh.has_text_frame:
                tf = sh.text_frame
                txt = tf.text.strip()
                # (3) 와이어 영역 문장형 텍스트
                if is_body and txt and (x + w) < DESC_X - 0.05:
                    is_marker = sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and txt.isdigit()
                    if not is_marker and len(txt.replace(" ", "")) > 4:
                        issues.append(f"[S{idx}] WIRE-TEXT '{txt[:20]}' (무언어 위반)")
                # (1) 오버플로
                ml = emu_in(tf.margin_left); mr = emu_in(tf.margin_right)
                mt = emu_in(tf.margin_top); mb = emu_in(tf.margin_bottom)
                aw = max(0.1, (w - ml - mr)) * 72
                th = 0.0; maxs = 0
                for p in tf.paragraphs:
                    if not p.runs:
                        continue
                    size = max((r.font.size.pt if r.font.size else 12) for r in p.runs)
                    maxs = max(maxs, size)
                    lt = "".join(r.text for r in p.runs)
                    lines = max(1, math.ceil(tw_pt(lt, size) / aw - 1e-6)) if tf.word_wrap else 1
                    ls = p.line_spacing or 1.0
                    sa = p.space_after.pt if p.space_after is not None else 0
                    th += lines * size * 1.2 * ls + sa
                ah = max(0.1, (h - mt - mb)) * 72
                if th > ah + maxs * 0.5:
                    issues.append(f"[S{idx}] TEXT-OVERFLOW '{_lbl(sh)}' need~{th/72:.2f} > {(h-mt-mb):.2f}in")
        if is_body and not has_frame:
            issues.append(f"[S{idx}] NO-FRAME (본문에 화면 전체 레이아웃 없음)")
    return issues


def _lbl(sh):
    if sh.has_text_frame:
        return sh.text_frame.text.replace("\n", " ")[:22]
    return str(sh.shape_type)


if __name__ == "__main__":
    p = sys.argv[1]
    iss = check(p)
    print(f"=== {p} : {len(iss)} findings ===")
    for i in iss:
        print(i)
