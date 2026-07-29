# -*- coding: utf-8 -*-
"""화면정의서 pptx 기하·규칙 감사 9축.

실행: python check_layout.py user-screens.pptx admin-screens.pptx   (0 findings = 통과)

감사 9축
  1) 경계          모든 도형이 콘텐츠 영역(MARGIN 기준 상자) 안에 있는가
  2) 마커 지름      빨간 번호 마커가 정원이며 규격 지름(1자리 0.26/0.24, 2자리 0.30/0.28)인가
  3) 마커 충돌      같은 슬라이드의 마커끼리 겹치지 않는가
  4) 마커 인접성    와이어프레임 마커가 가리키는 실루엣 요소 곁(0.30in 이내)에 놓였는가
  5) 마커·설명 번호 일치  와이어프레임 마커 번호 집합 = 우측 Description 마커 번호 집합인가
  6) 프레임 소속    와이어 마커는 브라우저 프레임 안에, 설명 마커는 Description 열 안에 있는가
  7) 톤다운 테두리  선 색이 흑·백·회색 팔레트 + 마커 빨강(#DC2626)으로만 구성되는가
  8) 행 정렬       동일 행으로 판정되는 도형군의 top이 일치하는가 (허용오차 ALIGN_TOL)
  9) 열 정렬       동일 열로 판정되는 도형군의 left가 일치하는가 (허용오차 ALIGN_TOL)

8·9축의 "동일 행/열" 판정 — 우연한 크기 일치와 의도적 오프셋을 배제하기 위해 아래를 모두 요구한다.
  · 대상: 네이티브 도형(AUTO_SHAPE)만. 텍스트 상자·커넥터·표는 제외한다
          (본문 텍스트와 실루엣 선은 좌표가 아니라 흐름으로 배치되므로 정렬 대상이 아니다).
  · 장식 제외: 흐름도 화살표(ARROW_SHAPES), 번호 마커, 최소 변 MIN_SIDE 미만의 점·닷.
  · 합동 요구: 두 도형의 폭·높이가 모두 SIZE_TOL 이내로 같아야 한다 — 반복 셀(카드·레인·
    스텝 노드)만 하나의 행/열로 본다. 컨테이너와 그 자식, 크기가 다른 이웃은 짝짓지 않는다.
  · 행 판정: 가로로 겹치지 않고(나란히), 세로로 OVERLAP_MIN 이상 겹치며, top 차이가
    ROW_BAND 이내. 열 판정은 축을 바꾼 대칭 조건.
  · ROW_BAND/COL_BAND는 허용오차가 아니라 **판정 근접 기준**이다. 이보다 크게 벌어진 좌표는
    의도적으로 다른 행/열에 놓인 것으로 보고 대조하지 않는다. 따라서 본 축이 잡는 결함은
    (ALIGN_TOL, ROW_BAND] 구간의 **근접 어긋남**이며, 그 이상은 배치 설계의 영역이다.
"""
import sys
import math
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE

EMU = 914400
SW, SH = 13.333, 7.5
MARGIN = 0.42
CONTENT_R = SW - MARGIN
CONTENT_B = SH - MARGIN

WIRE_W = 5.65
WIRE_R = MARGIN + WIRE_W
DESC_X = WIRE_R + 0.3

BOUND_TOL = 0.02          # 반올림 오차 허용
MARKER_D = (0.24, 0.26, 0.28, 0.30)
MARKER_D_TWO = 0.28       # 2자리 번호 최소 지름
ADJ_MAX = 0.30            # 마커–실루엣 요소 최대 이격 (실측 최악 0.12in)
LEAF_AREA_RATIO = 0.35    # 이 비율을 넘는 도형은 컨테이너로 보고 인접성 기준에서 제외
PALETTE = {"111111", "333333", "666666", "999999",
           "CCCCCC", "DDDDDD", "EEEEEE", "F7F7F7", "FFFFFF", "DC2626"}

ALIGN_TOL = 0.01          # 행/열 좌표 허용오차
ROW_BAND = 0.06           # 동일 행 판정: top 근접 기준
COL_BAND = 0.06           # 동일 열 판정: left 근접 기준
SIZE_TOL = 0.01           # 합동 판정: 폭·높이 일치 기준
OVERLAP_MIN = 0.6         # 동일 행/열 판정: 교차 축 겹침 최소 비율
MIN_SIDE = 0.12           # 이보다 작은 변의 도형은 장식(점·닷)으로 보고 제외
ARROW_SHAPES = {MSO_SHAPE.RIGHT_ARROW, MSO_SHAPE.LEFT_ARROW,
                MSO_SHAPE.UP_ARROW, MSO_SHAPE.DOWN_ARROW}


class Box:
    __slots__ = ("x", "y", "w", "h", "text", "marker", "kind")

    def __init__(self, x, y, w, h, text, marker, kind):
        self.x = x; self.y = y; self.w = w; self.h = h
        self.text = text; self.marker = marker; self.kind = kind

    @property
    def r(self):
        return self.x + self.w

    @property
    def b(self):
        return self.y + self.h

    @property
    def area(self):
        return self.w * self.h


def _collect(slide):
    boxes = []
    for sh in slide.shapes:
        try:
            x = sh.left / EMU; y = sh.top / EMU
            w = sh.width / EMU; h = sh.height / EMU
        except TypeError:
            continue
        text = sh.text_frame.text.strip() if sh.has_text_frame else ""
        is_marker = sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and text.isdigit()
        boxes.append((sh, Box(x, y, w, h, text, is_marker, sh.shape_type)))
    return boxes


def _edge_gap(a, b):
    dx = max(b.x - a.r, a.x - b.r, 0.0)
    dy = max(b.y - a.b, a.y - b.b, 0.0)
    return math.hypot(dx, dy)


def _overlap(a, b):
    return min(a.r, b.r) - max(a.x, b.x) > 0.0 and min(a.b, b.b) - max(a.y, b.y) > 0.0


def _line_rgb(sh):
    try:
        if sh.line.fill.type is None:
            return None
        return str(sh.line.color.rgb)
    except Exception:
        return None


def _label(bx):
    return (bx.text.replace("\n", " ")[:22] or str(bx.kind))


def _span_overlap(a0, a1, b0, b1):
    return min(a1, b1) - max(a0, b0)


def _align_targets(pairs):
    """8·9축 대상 도형. 텍스트 상자·커넥터·표·화살표·마커·미세 장식을 걸러낸다."""
    out = []
    for sh, b in pairs:
        if b.kind != MSO_SHAPE_TYPE.AUTO_SHAPE or b.marker:
            continue
        try:
            if sh.auto_shape_type in ARROW_SHAPES:
                continue
        except (AttributeError, ValueError):
            pass
        if b.w < MIN_SIDE or b.h < MIN_SIDE:
            continue
        out.append(b)
    return out


def _congruent(a, b):
    """반복 셀 여부 — 폭·높이가 모두 일치해야 하나의 행/열 구성원으로 본다."""
    return abs(a.w - b.w) <= SIZE_TOL and abs(a.h - b.h) <= SIZE_TOL


def _align_issues(idx, targets):
    """(8) 행 정렬 · (9) 열 정렬."""
    issues = []
    for i in range(len(targets)):
        for j in range(i + 1, len(targets)):
            a, b = targets[i], targets[j]
            if not _congruent(a, b):
                continue
            # (8) 행: 나란히 놓이고 세로로 겹치는 합동 도형 → top 일치
            if (_span_overlap(a.x, a.r, b.x, b.r) <= 0.0
                    and _span_overlap(a.y, a.b, b.y, b.b) >= OVERLAP_MIN * min(a.h, b.h)
                    and abs(a.y - b.y) <= ROW_BAND
                    and abs(a.y - b.y) > ALIGN_TOL):
                issues.append(f"[S{idx}] ROW-ALIGN '{_label(a)}'<->'{_label(b)}' "
                              f"top={a.y:.3f} vs {b.y:.3f} (Δ{abs(a.y - b.y):.3f}in "
                              f"> {ALIGN_TOL}in)")
            # (9) 열: 겹겹이 쌓이고 가로로 겹치는 합동 도형 → left 일치
            if (_span_overlap(a.y, a.b, b.y, b.b) <= 0.0
                    and _span_overlap(a.x, a.r, b.x, b.r) >= OVERLAP_MIN * min(a.w, b.w)
                    and abs(a.x - b.x) <= COL_BAND
                    and abs(a.x - b.x) > ALIGN_TOL):
                issues.append(f"[S{idx}] COL-ALIGN '{_label(a)}'<->'{_label(b)}' "
                              f"left={a.x:.3f} vs {b.x:.3f} (Δ{abs(a.x - b.x):.3f}in "
                              f"> {ALIGN_TOL}in)")
    return issues


def check(path):
    prs = Presentation(path)
    issues = []
    for idx, slide in enumerate(prs.slides, 1):
        pairs = _collect(slide)
        boxes = [b for _, b in pairs]
        is_body = any(sh.has_table for sh in slide.shapes)

        markers = [b for b in boxes if b.marker]
        wire_mk = [b for b in markers if b.x < DESC_X]
        desc_mk = [b for b in markers if b.x >= DESC_X - 0.05]

        # 브라우저 프레임 상단 = 본문 와이어 영역의 세로 시작점
        frame_top = None
        for b in boxes:
            if abs(b.x - MARGIN) < 0.15 and 5.3 < b.w < 5.9 and b.h > 3.0:
                frame_top = b.y if frame_top is None else min(frame_top, b.y)

        for sh, b in pairs:
            # (1) 경계
            if (b.x < MARGIN - BOUND_TOL or b.y < MARGIN - BOUND_TOL
                    or b.r > CONTENT_R + BOUND_TOL or b.b > CONTENT_B + BOUND_TOL):
                issues.append(f"[S{idx}] BOUNDS '{_label(b)}' "
                              f"l={b.x:.2f} t={b.y:.2f} r={b.r:.2f} b={b.b:.2f}")
            # (7) 톤다운 테두리
            rgb = _line_rgb(sh)
            if rgb is not None and rgb.upper() not in PALETTE:
                issues.append(f"[S{idx}] LINE-COLOR '{_label(b)}' #{rgb} (팔레트 외)")

        # (2) 마커 지름
        for m in markers:
            if abs(m.w - m.h) > 0.005:
                issues.append(f"[S{idx}] MARKER-SHAPE '{m.text}' {m.w:.3f}x{m.h:.3f} (정원 아님)")
            if not any(abs(m.w - d) < 0.005 for d in MARKER_D):
                issues.append(f"[S{idx}] MARKER-DIAM '{m.text}' d={m.w:.3f} (규격 외)")
            if len(m.text) >= 2 and m.w < MARKER_D_TWO - 0.005:
                issues.append(f"[S{idx}] MARKER-DIAM-2DIGIT '{m.text}' d={m.w:.3f} "
                              f"(2자리는 {MARKER_D_TWO} 이상)")

        # (3) 마커 충돌
        for i in range(len(markers)):
            for j in range(i + 1, len(markers)):
                if _overlap(markers[i], markers[j]):
                    issues.append(f"[S{idx}] MARKER-COLLISION '{markers[i].text}'"
                                  f"<->'{markers[j].text}'")

        # (4) 마커 인접성
        if wire_mk:
            inner_area = WIRE_W * max(0.1, CONTENT_B - (frame_top or MARGIN))
            leaves = [b for b in boxes
                      if not b.marker and b.x < DESC_X
                      and b.w > 0.03 and b.h > 0.03
                      and b.area <= inner_area * LEAF_AREA_RATIO]
            for m in wire_mk:
                if not leaves:
                    issues.append(f"[S{idx}] MARKER-ADJACENCY '{m.text}' (인접 실루엣 요소 없음)")
                    continue
                g = min(_edge_gap(m, o) for o in leaves)
                if g > ADJ_MAX:
                    issues.append(f"[S{idx}] MARKER-ADJACENCY '{m.text}' gap={g:.2f}in "
                                  f"> {ADJ_MAX}in")

        # (5) 마커·설명 번호 일치
        if is_body:
            wn = sorted(int(m.text) for m in wire_mk)
            dn = sorted(int(m.text) for m in desc_mk)
            if wn != dn:
                issues.append(f"[S{idx}] MARKER-DESC-MISMATCH wire={wn} desc={dn}")

        # (6) 프레임 소속
        if is_body and frame_top is not None:
            for m in wire_mk:
                if (m.x < MARGIN - BOUND_TOL or m.r > WIRE_R + BOUND_TOL
                        or m.y < frame_top - BOUND_TOL or m.b > CONTENT_B + BOUND_TOL):
                    issues.append(f"[S{idx}] FRAME-MEMBER '{m.text}' 와이어 마커가 프레임 밖")
            for m in desc_mk:
                if m.x < DESC_X - BOUND_TOL or m.r > CONTENT_R + BOUND_TOL:
                    issues.append(f"[S{idx}] FRAME-MEMBER '{m.text}' 설명 마커가 Description 열 밖")

        # (8) 행 정렬 · (9) 열 정렬
        issues.extend(_align_issues(idx, _align_targets(pairs)))
    return issues


if __name__ == "__main__":
    targets = sys.argv[1:] or ["user-screens.pptx", "admin-screens.pptx"]
    total = 0
    for p in targets:
        found = check(p)
        total += len(found)
        print(f"=== {p} : {len(found)} findings ===")
        for f in found:
            print(f)
    sys.exit(1 if total else 0)
