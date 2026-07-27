# -*- coding: utf-8 -*-
"""
AX Platform 화면정의서 v3 공용 헬퍼 (PHASE 1·2 공유) — 확정 템플릿
- 본문: 무언어 무채색 실루엣 와이어프레임 + 빨간 마커 → 우측 Description
        + 하단 "화면 룰·기획 근거" 요약 박스(전 화면 표준: 설계 배경 3~5줄 + BE·DB 한 줄)
- p0: 좌[상 구조 트리 / 하 사용자 흐름 + 설계 근거 주석] / 우[정의·역할·목적·기획 의도·룰·개발 연동 노트 5단]
- 네이티브 도형/텍스트/표만. 색은 흑·백·회색 + 빨간 마커(#DC2626)만 예외. 맑은 고딕. 16:9.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

FONT = "맑은 고딕"
AUTHOR = ""   # 빌드 스크립트에서 spec_common.AUTHOR = "..." 로 지정 (헤더 표 '작성자' 칸)

INK   = RGBColor(0x11, 0x11, 0x11)
G33   = RGBColor(0x33, 0x33, 0x33)
G66   = RGBColor(0x66, 0x66, 0x66)
G99   = RGBColor(0x99, 0x99, 0x99)
GCC   = RGBColor(0xCC, 0xCC, 0xCC)
GDD   = RGBColor(0xDD, 0xDD, 0xDD)
GEE   = RGBColor(0xEE, 0xEE, 0xEE)
GF7   = RGBColor(0xF7, 0xF7, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED   = RGBColor(0xDC, 0x26, 0x26)

SW, SH = 13.333, 7.5
MARGIN = 0.42
CONTENT_X = MARGIN
CONTENT_W = SW - 2 * MARGIN
CONTENT_BOTTOM = SH - MARGIN
EMU = 914400

_TEMPLATE = os.environ.get(
    "PPTX_TEMPLATE",
    r"C:\Users\USER\AppData\Local\Temp\claude\i--Workspace\d5b0a822-1399-4b83-8fda-ea38b58409d1\scratchpad\default_template.pptx",
)


def new_deck():
    prs = Presentation(_TEMPLATE) if os.path.exists(_TEMPLATE) else Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)
    return prs


def add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _fill(sp, color):
    if color is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = color


def _line(sp, color, w):
    if color is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = color; sp.line.width = Pt(w)


def rect(slide, x, y, w, h, fill=WHITE, line=GCC, line_w=1.5, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    _fill(sp, fill); _line(sp, line, line_w); sp.shadow.inherit = False
    return sp


def seg(slide, x1, y1, x2, y2, color=G99, w=1.2):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cn.line.color.rgb = color; cn.line.width = Pt(w); cn.shadow.inherit = False
    return cn


def _apply_par(p, runs, align, space_after, ls):
    p.alignment = align
    if space_after is not None:
        p.space_after = Pt(space_after)
    p.space_before = Pt(0)
    if ls is not None:
        p.line_spacing = ls
    for (text, size, bold, color) in runs:
        r = p.add_run(); r.text = text
        r.font.name = FONT; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
        rPr = r._r.get_or_add_rPr()
        ea = rPr.makeelement(qn('a:ea'), {}); ea.set('typeface', FONT); rPr.append(ea)


def textbox(slide, x, y, w, h, paragraphs, anchor=MSO_ANCHOR.TOP, wrap=True,
            left_in=0.06, right_in=0.06, top_in=0.03):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = Inches(left_in); tf.margin_right = Inches(right_in)
    tf.margin_top = Inches(top_in); tf.margin_bottom = Inches(0.02)
    tf.auto_size = MSO_AUTO_SIZE.NONE
    for i, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _apply_par(p, para.get("runs", []), para.get("align", PP_ALIGN.LEFT),
                   para.get("space_after", 3), para.get("line_spacing", 1.0))
    return tb


def line_para(text, size=10.5, bold=False, color=G33, align=PP_ALIGN.LEFT, space_after=3, ls=1.05):
    return {"runs": [(text, size, bold, color)], "align": align, "space_after": space_after, "line_spacing": ls}


# ============================================================
# 헤더 표
# ============================================================
def _cell_border(cell, color="666666", w=9525):
    tcPr = cell._tc.get_or_add_tcPr()
    for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        for ex in tcPr.findall(qn(tag)):
            tcPr.remove(ex)
        ln = tcPr.makeelement(qn(tag), {"w": str(w), "cap": "flat"})
        fill = ln.makeelement(qn("a:solidFill"), {})
        clr = fill.makeelement(qn("a:srgbClr"), {"val": color})
        fill.append(clr); ln.append(fill); tcPr.append(ln)


def _cell(cell, text, size=10, bold=False, color=INK, fill=WHITE, align=PP_ALIGN.LEFT):
    cell.fill.solid(); cell.fill.fore_color.rgb = fill
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.06)
    cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = FONT; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    rPr = r._r.get_or_add_rPr(); ea = rPr.makeelement(qn('a:ea'), {}); ea.set('typeface', FONT); rPr.append(ea)
    _cell_border(cell)


def header_table(slide, screen_name, title, screen_id, path, page_no, page_total, subtitle=""):
    x, y, w = CONTENT_X, MARGIN, CONTENT_W
    h = 1.14
    g = slide.shapes.add_table(3, 4, Inches(x), Inches(y), Inches(w), Inches(h)).table
    g.first_row = False; g.horz_banding = False
    cw = [1.5, 4.7, 1.4, w - 1.5 - 4.7 - 1.4]
    for i, c in enumerate(cw):
        g.columns[i].width = Inches(c)
    g.rows[0].height = Inches(0.42); g.rows[1].height = Inches(0.36); g.rows[2].height = Inches(0.36)
    a = g.cell(0, 0); a.merge(g.cell(0, 2))
    head = screen_name + (f"  —  {subtitle}" if subtitle else "")
    _cell(a, head, size=13, bold=True, color=WHITE, fill=INK)
    _cell(g.cell(0, 3), f"Page {page_no} / {page_total}", size=10.5, bold=True, color=WHITE, fill=INK, align=PP_ALIGN.CENTER)
    _cell(g.cell(1, 0), "Page Title", size=9.5, bold=True, color=G33, fill=GEE)
    _cell(g.cell(1, 1), title, size=10.5, color=INK)
    _cell(g.cell(1, 2), "Screen ID", size=9.5, bold=True, color=G33, fill=GEE)
    _cell(g.cell(1, 3), screen_id, size=10.5, bold=True, color=INK)
    _cell(g.cell(2, 0), "Screen Path", size=9.5, bold=True, color=G33, fill=GEE)
    _cell(g.cell(2, 1), path, size=9.5, color=G33)
    _cell(g.cell(2, 2), "작성자", size=9.5, bold=True, color=G33, fill=GEE)
    _cell(g.cell(2, 3), AUTHOR, size=10)
    return y + h


# ============================================================
# 마커
# ============================================================
def marker(slide, x, y, num, d=0.26):
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    _fill(sp, RED); _line(sp, WHITE, 1.0); sp.shadow.inherit = False
    tf = sp.text_frame
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(num)
    r.font.name = FONT; r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = WHITE
    return sp


# ============================================================
# 실루엣 (무채색). active → #333 계열 / dim → #DDD
# ============================================================
def _oc(active):
    return G33 if active else GDD


def sil_box(s, x, y, w, h, active, lw=1.5, fill=WHITE, round=False):
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if round else MSO_SHAPE.RECTANGLE
    return rect(s, x, y, w, h, fill=fill, line=_oc(active), line_w=lw, shape=shp)


def sil_button(s, x, y, w, h, active):
    rect(s, x, y, w, h, fill=(GF7 if active else WHITE), line=_oc(active), line_w=1.2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)


def sil_input(s, x, y, w, h, active):
    rect(s, x, y, w, h, fill=WHITE, line=(GCC if active else GEE), line_w=1.0)
    seg(s, x, y + h, x + w, y + h, color=_oc(active), w=1.8)


def sil_lines(s, x, y, w, active, n=3, gap=0.135, fracs=None):
    col = G66 if active else GDD
    fr = fracs or ([1.0] * (n - 1) + [0.62])
    for i in range(n):
        yy = y + i * gap
        seg(s, x, yy, x + w * fr[i], yy, color=col, w=1.1)


def sil_grid(s, x, y, w, h, active, cols, rows, gap=0.12):
    cw = (w - gap * (cols - 1)) / cols
    ch = (h - gap * (rows - 1)) / rows
    for r in range(rows):
        for c in range(cols):
            rect(s, x + c * (cw + gap), y + r * (ch + gap), cw, ch,
                 fill=WHITE, line=_oc(active), line_w=1.2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)


def sil_bars(s, x, y, w, h, active, pattern=None):
    pat = pattern or [0.55, 0.8, 0.4, 0.95, 0.6, 0.72, 0.5]
    n = len(pat)
    seg(s, x, y + h, x + w, y + h, color=_oc(active), w=1.3)  # 축
    bw = (w / n) * 0.6
    step = w / n
    fill = GCC if active else GEE
    for i, p in enumerate(pat):
        bh = h * p
        bx = x + i * step + (step - bw) / 2
        rect(s, bx, y + h - bh, bw, bh, fill=fill, line=(_oc(active)), line_w=0.75)


def sil_linechart(s, x, y, w, h, active, pts=None):
    p = pts or [0.7, 0.45, 0.6, 0.3, 0.5, 0.2, 0.35]
    n = len(p)
    seg(s, x, y + h, x + w, y + h, color=_oc(active), w=1.3)
    col = G66 if active else GDD
    step = w / (n - 1)
    for i in range(n - 1):
        seg(s, x + i * step, y + h * p[i], x + (i + 1) * step, y + h * p[i + 1], color=col, w=1.5)


def sil_pill(s, x, y, w, h, active):
    rect(s, x, y, w, h, fill=WHITE, line=_oc(active), line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)


def sil_circle(s, x, y, d, active):
    sp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    _fill(sp, WHITE); _line(sp, _oc(active), 1.2); sp.shadow.inherit = False


def tiny(s, x, y, w, text, active, align=PP_ALIGN.CENTER, size=9):
    col = G99 if active else GDD
    textbox(s, x, y - 0.03, w, 0.26, [line_para(text, size=size, bold=False, color=col, align=align, space_after=0)],
            anchor=MSO_ANCHOR.MIDDLE, top_in=0.0)


# ============================================================
# 본문 슬라이드 프레임 + Description
# ============================================================
class Ctx:
    def __init__(self, active_set):
        self.active = active_set
    def on(self, n):
        return n in self.active
    def mk(self, s, n, x, y):
        if n in self.active:
            marker(s, x, y, n)


def body_frame(prs, screen_name, title, screen_id, path, page_no, page_total, subtitle=""):
    s = add_slide(prs)
    hb = header_table(s, screen_name, title, screen_id, path, page_no, page_total, subtitle)
    top = hb + 0.16
    wx, ww = CONTENT_X, 5.65
    # 브라우저 프레임
    rect(s, wx, top, ww, CONTENT_BOTTOM - top, fill=WHITE, line=G33, line_w=2.4)
    bar_h = 0.16
    rect(s, wx, top, ww, bar_h, fill=GF7, line=G33, line_w=2.4)
    for i, cx in enumerate((wx + 0.11, wx + 0.22, wx + 0.33)):
        sp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx), Inches(top + 0.055), Inches(0.05), Inches(0.05))
        _fill(sp, GCC); _line(sp, None, 0); sp.shadow.inherit = False
    pad = 0.14
    inner = {"x": wx + pad, "y": top + bar_h + pad, "w": ww - 2 * pad,
             "h": CONTENT_BOTTOM - (top + bar_h + pad) - pad}
    dx = wx + ww + 0.3
    return {"slide": s, "inner": inner,
            "desc_x": dx, "desc_w": CONTENT_X + CONTENT_W - dx,
            "desc_y": top, "desc_h": CONTENT_BOTTOM - top}


def _rule_box(s, x, y, w, box, max_h):
    """하단 '화면 룰·기획 근거' 요약 박스. box=(설계배경 줄 리스트, BE·DB 한 줄)."""
    bg_lines, be_line = box
    rect(s, x, y, w, max_h, fill=GF7, line=G66, line_w=1.0)
    seg(s, x, y + 0.30, x + w, y + 0.30, color=GCC, w=1.0)
    textbox(s, x + 0.12, y + 0.045, w - 0.24, 0.26,
            [line_para("화면 룰·기획 근거", size=9.5, bold=True, color=INK, space_after=0)])
    paras = []
    for b in bg_lines:
        paras.append({"runs": [("· ", 8.5, True, G99), (b, 8.5, False, G33)],
                      "align": PP_ALIGN.LEFT, "space_after": 1, "line_spacing": 1.04})
    if be_line:
        paras.append({"runs": [("BE·DB  ", 8, True, RED), (be_line, 8.5, False, G66)],
                      "align": PP_ALIGN.LEFT, "space_after": 0, "line_spacing": 1.04})
    textbox(s, x + 0.14, y + 0.34, w - 0.28, max_h - 0.38, paras, anchor=MSO_ANCHOR.TOP)


def desc_panel(area, items, rule_box=None, title="Description"):
    s = area["slide"]; x = area["desc_x"]; w = area["desc_w"]
    y = area["desc_y"]; h = area["desc_h"]
    textbox(s, x, y, w, 0.28, [line_para(title, size=11, bold=True, color=INK, space_after=0)])
    top = y + 0.36
    # 하단 규칙 박스 높이 산정(줄바꿈 반영)
    box_h = 0.0
    if rule_box:
        bg_lines, be_line = rule_box
        line_ct = sum(_wrap_lines(b, 8.5, w - 0.46) for b in bg_lines)
        if be_line:
            line_ct += _wrap_lines("BE·DB  " + be_line, 8.5, w - 0.46)
        box_h = 0.34 + line_ct * 0.182 + 0.12
        box_h = min(box_h, h * 0.5)
    avail = (y + h) - top - (box_h + 0.16 if rule_box else 0.0)
    n = len(items)
    row_h = avail / n if n else avail
    for i, it in enumerate(items):
        num, name, detail = it[0], it[1], it[2]
        rule = it[3] if len(it) > 3 else None
        ry = top + i * row_h
        if i > 0:
            seg(s, x, ry - 0.02, x + w, ry - 0.02, color=GEE, w=0.75)
        marker(s, x, ry + 0.03, num, d=0.24)
        runs = [(f"{name} ", 9.5, True, INK), ("— ", 9.5, False, G99), (detail, 9, False, G66)]
        paras = [{"runs": runs, "align": PP_ALIGN.LEFT, "space_after": 1, "line_spacing": 1.04}]
        if rule:
            paras.append({"runs": [("안내 ", 8, True, G99), (rule, 8.5, False, G33)],
                          "align": PP_ALIGN.LEFT, "space_after": 0, "line_spacing": 1.03})
        textbox(s, x + 0.34, ry, w - 0.34, row_h, paras, anchor=MSO_ANCHOR.TOP)
    if rule_box:
        _rule_box(s, x, top + avail + 0.16, w, rule_box, box_h)


def screen(prs, hdr, draw_fn, desc_items, active=None, page_no=1, page_total=1, subtitle="", rule_box=None):
    """hdr=(screen_name,title,screen_id,path). draw_fn(slide,inner,ctx). rule_box=(배경줄들, BE줄)."""
    area = body_frame(prs, hdr[0], hdr[1], hdr[2], hdr[3], page_no, page_total, subtitle)
    all_nums = set(it[0] for it in desc_items)
    act = all_nums if active is None else set(active)
    ctx = Ctx(act)
    draw_fn(area["slide"], area["inner"], ctx)
    desc_panel(area, [it for it in desc_items if it[0] in act], rule_box=rule_box)
    return area


# ============================================================
# p0 정의 슬라이드 (구조 트리 + 사용자 흐름 / 톤다운 4단)
# ============================================================
def _node_box(s, x, y, w, h, label, size=9.5, bold=False):
    rect(s, x, y, w, h, fill=WHITE, line=G66, line_w=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    textbox(s, x + 0.05, y, w - 0.1, h, [line_para(label, size=size, bold=bold, color=G33,
            align=PP_ALIGN.LEFT, space_after=0, ls=1.0)], anchor=MSO_ANCHOR.MIDDLE)


def p0_tree(s, x, y, w, flat):
    """flat: list of (depth, label). depth 0=root."""
    row_h = 0.30; gap = 0.075; indent = 0.42
    nh = 0.255
    last_cy = {}
    cy = y
    for (depth, label) in flat:
        nx = x + depth * indent
        nw = w - depth * indent
        if nw > 3.7:
            nw = max(2.2, nw * 0.62)
        bold = depth == 0
        _node_box(s, nx, cy, nw, nh, label, size=(10 if bold else 9), bold=bold)
        if depth > 0:
            px = x + (depth - 1) * indent + 0.12
            parent_cy = last_cy.get(depth - 1, cy)
            seg(s, px, parent_cy, px, cy + nh / 2, color=G99, w=1.3)          # 세로
            seg(s, px, cy + nh / 2, nx, cy + nh / 2, color=G99, w=1.3)        # 가로 stub
        last_cy[depth] = cy + nh / 2
        cy += row_h + gap
    return cy


def p0_flow(s, x, y, w, steps, branch=None):
    """steps: list of str (좌→우). branch=(after_index, label)."""
    n = len(steps)
    gap = 0.26
    bw = (w - gap * (n - 1)) / n
    bh = 0.5
    cx = x
    centers = []
    for i, st in enumerate(steps):
        _node_box(s, cx, y, bw, bh, st, size=8.5)
        centers.append((cx + bw / 2, y + bh / 2, cx + bw))
        if i < n - 1:
            ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(cx + bw + 0.02),
                                    Inches(y + bh / 2 - 0.07), Inches(gap - 0.04), Inches(0.14))
            _fill(ar, GCC); _line(ar, None, 0); ar.shadow.inherit = False
        cx += bw + gap
    if branch:
        idx, blabel = branch
        c = centers[idx]
        by = y + bh + 0.34
        da = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(c[0] - 0.07), Inches(y + bh + 0.02),
                                Inches(0.14), Inches(0.28))
        _fill(da, GCC); _line(da, None, 0); da.shadow.inherit = False
        _node_box(s, c[0] - bw / 2, by, bw, 0.42, blabel, size=8.5)


def p0_slots_flow(s, x, y, w, start, slots, merge, reject):
    """병렬 2-슬롯 분기·합류 흐름도: 신청 →[관계사/전사 병렬]→ 게시, 슬롯 영역에서 반려→종결.
    총 높이 ≈ 1.56in (block 0.96 + 반려 0.60)."""
    gap = 0.30
    nw = (w - 2 * gap) / 3
    lane_h = 0.40; lane_gap = 0.16
    block_h = lane_h * 2 + lane_gap
    c1 = x; c2 = x + nw + gap; c3 = c2 + nw + gap
    top_y = y; bot_y = y + lane_h + lane_gap
    scy = y + block_h / 2
    tcy = top_y + lane_h / 2; bcy = bot_y + lane_h / 2
    _node_box(s, c1, scy - 0.22, nw, 0.44, start, size=8.5)
    _node_box(s, c2, top_y, nw, lane_h, slots[0], size=8)
    _node_box(s, c2, bot_y, nw, lane_h, slots[1], size=8)
    _node_box(s, c3, scy - 0.22, nw, 0.44, merge, size=8.5, bold=True)
    # 분기 버스 (신청 → 두 슬롯)
    bx1 = (c1 + nw + c2) / 2
    seg(s, c1 + nw, scy, bx1, scy, color=G99, w=1.3)
    seg(s, bx1, tcy, bx1, bcy, color=G99, w=1.3)
    seg(s, bx1, tcy, c2, tcy, color=G99, w=1.3)
    seg(s, bx1, bcy, c2, bcy, color=G99, w=1.3)
    # 합류 버스 (두 슬롯 → 게시)
    bx2 = (c2 + nw + c3) / 2
    seg(s, c2 + nw, tcy, bx2, tcy, color=G99, w=1.3)
    seg(s, c2 + nw, bcy, bx2, bcy, color=G99, w=1.3)
    seg(s, bx2, tcy, bx2, bcy, color=G99, w=1.3)
    seg(s, bx2, scy, c3, scy, color=G99, w=1.3)
    ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(c3 - 0.15), Inches(scy - 0.05), Inches(0.13), Inches(0.11))
    _fill(ar, GCC); _line(ar, None, 0); ar.shadow.inherit = False
    # 반려 → 종결 (슬롯 영역 아래로 분기)
    da = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(c2 + nw / 2 - 0.07), Inches(y + block_h + 0.01),
                            Inches(0.14), Inches(0.22))
    _fill(da, GCC); _line(da, None, 0); da.shadow.inherit = False
    _node_box(s, c2, y + block_h + 0.24, nw, 0.36, reject, size=8)


def _wrap_lines(text, size, box_w_in):
    """맑은 고딕 근사 줄 수 추정 (와이드 문자=size, ASCII≈0.55*size)."""
    import unicodedata
    wpt = 0.0
    for ch in text:
        wpt += size * (1.0 if unicodedata.east_asian_width(ch) in ("W", "F", "A") else 0.55)
    avail = max(0.1, box_w_in) * 72
    import math
    return max(1, math.ceil(wpt / avail - 1e-6))


def p0_side(s, x, y, w, h, sections):
    """sections: list of (title, [bullets]). 내용 분량(줄바꿈 반영) 맞춰 세로 배치."""
    line_h = 0.205; head = 0.30; pad = 0.13
    inner_w = w - 0.28 - 0.15  # 불렛 '· ' 들여쓰기 여유
    heights = []
    for (_, bullets) in sections:
        lines = sum(_wrap_lines(b, 9, inner_w) for b in bullets)
        heights.append(head + lines * line_h + pad)
    total = sum(heights)
    gap = 0.14
    avail = h - gap * (len(sections) - 1)
    factor = avail / total   # 세로 여백 없이 컬럼 채움
    cy = y
    for (title, bullets), hh0 in zip(sections, heights):
        hh = hh0 * factor
        rect(s, x, cy, w, hh, fill=WHITE, line=GCC, line_w=1.0)
        textbox(s, x + 0.12, cy + 0.04, w - 0.24, 0.26,
                [line_para(title, size=11, bold=True, color=G33, space_after=0)])
        seg(s, x + 0.12, cy + 0.30, x + w - 0.12, cy + 0.30, color=GEE, w=1.0)
        paras = []
        for b in bullets:
            paras.append({"runs": [("· ", 9, True, G99), (b, 9, False, G66)],
                          "align": PP_ALIGN.LEFT, "space_after": 2, "line_spacing": 1.05})
        textbox(s, x + 0.14, cy + 0.34, w - 0.28, hh - 0.38, paras, anchor=MSO_ANCHOR.TOP)
        cy += hh + gap


def _flow_note(s, x, y, w, notes):
    """사용자 흐름 아래 설계 근거 주석 2~4줄."""
    textbox(s, x, y, w, 0.22, [line_para("설계 근거", size=9, bold=True, color=G99, space_after=0)])
    paras = []
    for t in notes:
        paras.append({"runs": [("· ", 8.5, True, G99), (t, 8.5, False, G66)],
                      "align": PP_ALIGN.LEFT, "space_after": 1, "line_spacing": 1.05})
    textbox(s, x + 0.02, y + 0.24, w - 0.04, 0.24 + len(notes) * 0.22, paras, anchor=MSO_ANCHOR.TOP)


def def_slide(prs, screen_id, screen_name, tree, flow, sections, flow_branch=None,
              flow_note=None, flow_slots=None, tree_caption="화면 구조", flow_caption="사용자 흐름"):
    s = add_slide(prs)
    # 타이틀 행
    textbox(s, CONTENT_X, MARGIN, CONTENT_W * 0.7, 0.4,
            [line_para(f"{screen_id}  ·  {screen_name}", size=15, bold=True, color=INK, space_after=0)],
            anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, CONTENT_X + CONTENT_W * 0.6, MARGIN, CONTENT_W * 0.4, 0.4,
            [line_para("화면 정의 (Screen Definition)", size=10, color=G99, align=PP_ALIGN.RIGHT, space_after=0)],
            anchor=MSO_ANCHOR.MIDDLE)
    seg(s, CONTENT_X, MARGIN + 0.44, CONTENT_X + CONTENT_W, MARGIN + 0.44, color=GCC, w=1.5)

    top = MARGIN + 0.58
    left_x = CONTENT_X
    left_w = CONTENT_W * 0.55
    right_x = CONTENT_X + left_w + 0.3
    right_w = CONTENT_X + CONTENT_W - right_x

    # 좌측: [상] 구조 트리
    textbox(s, left_x, top, left_w, 0.24, [line_para(tree_caption, size=9.5, bold=True, color=G99, space_after=0)])
    tree_bottom = p0_tree(s, left_x + 0.05, top + 0.3, left_w - 0.1, tree)

    # 좌측: [하] 사용자 흐름 + 설계 근거 주석 — 하단 기준으로 배치
    notes = flow_note or []
    note_h = (0.26 + len(notes) * 0.22) if notes else 0.0
    if flow_slots:
        flow_block = 1.56                                     # 병렬 2-슬롯 분기·합류 흐름도
    else:
        flow_block = 0.5 + (0.78 if flow_branch else 0.0)     # 노드 + (분기 갈래)
    needed = 0.24 + 0.06 + flow_block + (0.14 + note_h if notes else 0.0)
    flow_cap_y = CONTENT_BOTTOM - needed
    if flow_cap_y < tree_bottom + 0.14:
        flow_cap_y = tree_bottom + 0.14
    textbox(s, left_x, flow_cap_y, left_w, 0.24, [line_para(flow_caption, size=9.5, bold=True, color=G99, space_after=0)])
    if flow_slots:
        p0_slots_flow(s, left_x + 0.05, flow_cap_y + 0.3, left_w - 0.1, *flow_slots)
    else:
        p0_flow(s, left_x + 0.05, flow_cap_y + 0.3, left_w - 0.1, flow, branch=flow_branch)
    if notes:
        _flow_note(s, left_x + 0.05, flow_cap_y + 0.3 + flow_block + 0.14, left_w - 0.1, notes)

    # 우측: 톤다운 5단 (정의·역할 / 목적 / 기획 의도 / 지켜야 할 룰 / 개발 연동 노트)
    p0_side(s, right_x, top, right_w, CONTENT_BOTTOM - top, sections)
    return s
